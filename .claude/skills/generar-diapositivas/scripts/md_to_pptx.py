#!/usr/bin/env python3
"""Genera un .pptx (y su .pdf) a partir de una nota de teoría de unidad,
siguiendo la convención de notas_unidades/**/*.md descrita en el skill
`notas-unidad`. Objetivo: el .pptx nunca "drifta" del .md porque se
regenera completo desde el .md cada vez; no se edita el .pptx a mano.

Uso:
    python md_to_pptx.py <ruta.md> [<ruta_salida.pptx>]

Si no se da ruta de salida, se usa el mismo nombre que el .md con
extensión .pptx, en el mismo directorio.
"""
import io
import re
import sys
import os
import math
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from PIL import ImageFont

# --------------------------------------------------------------------------
# Geometría y paleta (heredadas del diseño original de la unidad 1)
# --------------------------------------------------------------------------

EMU_PER_PT = 12700
EMU_PER_IN = 914400

# Marcador que excluye una sección de las diapositivas (queda solo en el .md).
SKIP_SLIDE_MARKER = "<!-- diapositivas: omitir -->"

SLIDE_W = 12192000
SLIDE_H = 6858000

MARGIN_L = 640080
MARGIN_R = 640080
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # 10911840

EYEBROW_TOP = 457200
EYEBROW_H = 320040
TITLE_TOP = 749808
FOOTER_TOP = 6400800
PAGENUM_TOP = 6446520
BODY_BOTTOM_MAX = 6190488  # deja espacio a la línea de fuente/página

NAVY = RGBColor(0x1E, 0x27, 0x61)
BODY_COLOR = RGBColor(0x27, 0x31, 0x4D)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GRAY = RGBColor(0x5B, 0x64, 0x82)
LGRAY = RGBColor(0xA6, 0xAD, 0xC7)
CARD_BG = RGBColor(0xF2, 0xF4, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEAD_BG = RGBColor(0x1E, 0x27, 0x61)
TABLE_ROW_ALT = RGBColor(0xF2, 0xF4, 0xFA)

FONT_TITLE = "Cambria"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

FONT_FILES = {
    ("Calibri", False, False): "calibri.ttf",
    ("Calibri", True, False): "calibrib.ttf",
    ("Calibri", False, True): "calibrii.ttf",
    ("Calibri", True, True): "calibriz.ttf",
    ("Cambria", False, False): "cambria.ttc",
    ("Cambria", True, False): "cambriab.ttf",
}
_FONT_CACHE = {}


def get_font(name, bold=False, italic=False, size_pt=12):
    key = (name, bold, italic, size_pt)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    fname = FONT_FILES.get((name, bold, italic), "calibri.ttf")
    path = os.path.join("C:/Windows/Fonts", fname)
    try:
        font = ImageFont.truetype(path, size_pt)
    except Exception:
        font = ImageFont.truetype("C:/Windows/Fonts/calibri.ttf", size_pt)
    _FONT_CACHE[key] = font
    return font


def text_width_px(text, name=FONT_BODY, bold=False, italic=False, size_pt=12):
    font = get_font(name, bold, italic, size_pt)
    bbox = font.getbbox(text)
    return bbox[2] - bbox[0]


def emu_to_pt72(emu):
    return emu / EMU_PER_IN * 72.0


def wrap_runs_to_lines(runs, max_width_pt, size_pt):
    """runs: list of (text,bold,italic,code). Devuelve lista de líneas,
    cada línea es lista de runs (mismo formato) que caben en max_width_pt."""
    words = []
    for (t, b, i, c) in runs:
        for j, w in enumerate(t.split(" ")):
            if w == "" and j == 0:
                continue
            words.append((w, b, i, c))
            words.append((" ", b, i, c))
    lines = []
    cur = []
    cur_w = 0.0
    for (w, b, i, c) in words:
        fname = FONT_CODE if c else FONT_TITLE if False else FONT_BODY
        ww = text_width_px(w, fname, b, i, size_pt)
        if cur and cur_w + ww > max_width_pt:
            lines.append(cur)
            cur = []
            cur_w = 0.0
        if w == " " and not cur:
            continue
        cur.append((w, b, i, c))
        cur_w += ww
    if cur:
        lines.append(cur)
    return lines


def count_wrapped_lines(text, max_width_pt, size_pt, bold=False, italic=False, font=FONT_BODY):
    if not text:
        return 1
    w = text_width_px(text, font, bold, italic, size_pt)
    if w <= max_width_pt:
        return 1
    return max(1, math.ceil(w / max_width_pt))


# --------------------------------------------------------------------------
# Utilidades de números romanos
# --------------------------------------------------------------------------

def to_roman(n):
    vals = [(10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I")]
    out = ""
    for v, s in vals:
        while n >= v:
            out += s
            n -= v
    return out


# --------------------------------------------------------------------------
# Conversión de LaTeX simple a texto unicode (para matemática en línea)
# --------------------------------------------------------------------------

SUPER_MAP = {"0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴", "5": "⁵",
             "6": "⁶", "7": "⁷", "8": "⁸", "9": "⁹", "+": "⁺", "-": "⁻",
             "N": "ᴺ", "n": "ⁿ", "t": "ᵗ", "r": "ʳ", "i": "ⁱ"}
SUB_MAP = {"0": "₀", "1": "₁", "2": "₂", "3": "₃", "4": "₄", "5": "₅",
           "6": "₆", "7": "₇", "8": "₈", "9": "₉", "t": "ₜ", "n": "ₙ"}


def _latex_frag_to_text(s):
    s = re.sub(r"\\d?frac\{([^{}]+)\}\{([^{}]+)\}", r"(\1)/(\2)", s)
    s = s.replace(r"\times", "×").replace(r"\neq", "≠").replace(r"\cdot", "·")
    s = s.replace(r"\dots", "…").replace(r"\ldots", "…")
    s = s.replace(r"\qquad", "   ").replace(r"\quad", " ")
    s = s.replace(r"\approx", "≈").replace(r"\leq", "≤").replace(r"\geq", "≥")

    def sup_repl(m):
        body = m.group(1) if m.group(1) is not None else m.group(2)
        return "".join(SUPER_MAP.get(c, c) for c in body)

    s = re.sub(r"\^\{([^{}]+)\}|\^(\S)", sup_repl, s)

    def sub_repl(m):
        body = m.group(1) if m.group(1) is not None else m.group(2)
        return "".join(SUB_MAP.get(c, c) for c in body)

    s = re.sub(r"_\{([^{}]+)\}|_(\S)", sub_repl, s)
    s = s.replace("\\", "")
    return s


def convert_inline_math(text):
    text = text.replace(r"\$", "\uE000")

    def repl(m):
        return _latex_frag_to_text(m.group(1))

    text = re.sub(r"\$([^$]+)\$", repl, text)
    text = text.replace("\uE000", "$")
    return text


LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")


def strip_links(text):
    return LINK_RE.sub(r"\1", text)


INLINE_TOKEN_RE = re.compile(r"(\*\*.+?\*\*|`[^`]+?`|\*[^*\n]+?\*)")


def parse_inline(text):
    """Devuelve lista de runs (text,bold,italic,code)."""
    text = convert_inline_math(text)
    text = strip_links(text)
    tokens = []
    pos = 0
    for m in INLINE_TOKEN_RE.finditer(text):
        if m.start() > pos:
            tokens.append((text[pos:m.start()], False, False, False))
        s = m.group(0)
        if s.startswith("**"):
            tokens.append((s[2:-2], True, False, False))
        elif s.startswith("`"):
            tokens.append((s[1:-1], False, False, True))
        else:
            tokens.append((s[1:-1], False, True, False))
        pos = m.end()
    if pos < len(text):
        tokens.append((text[pos:], False, False, False))
    return [t for t in tokens if t[0] != ""]


def plain_text_of(text):
    return "".join(t for t, b, i, c in parse_inline(text))


# --------------------------------------------------------------------------
# Parser de markdown -> estructura de documento
# --------------------------------------------------------------------------

def parse_table(lines):
    rows = []
    for ln in lines:
        s = ln.strip()
        if re.match(r"^\|[\s:|-]+\|$", s):
            continue
        cells = [c.strip() for c in s.strip("|").split("|")]
        rows.append(cells)
    return rows


def parse_blockquote(bq_lines):
    blocks = []
    i, n = 0, len(bq_lines)
    while i < n:
        if bq_lines[i] == "":
            i += 1
            continue
        if re.match(r"^[-*]\s+", bq_lines[i]):
            items = []
            while i < n and re.match(r"^[-*]\s+", bq_lines[i]):
                items.append((re.sub(r"^[-*]\s+", "", bq_lines[i]), 0))
                i += 1
            blocks.append({"type": "bullets", "items": items})
            continue
        para = [bq_lines[i]]
        i += 1
        while i < n and bq_lines[i] != "" and not re.match(r"^[-*]\s+", bq_lines[i]):
            para.append(bq_lines[i])
            i += 1
        blocks.append({"type": "para", "text": " ".join(para)})
    return blocks


def consume_skip_marker(lines, i):
    """Si la próxima línea no vacía es SKIP_SLIDE_MARKER, la consume y
    devuelve (True, nuevo_i); si no, (False, i sin modificar)."""
    n = len(lines)
    j = i
    while j < n and lines[j].strip() == "":
        j += 1
    if j < n and lines[j].strip() == SKIP_SLIDE_MARKER:
        return True, j + 1
    return False, i


def parse_blocks_until_heading(lines, i, stop_levels=("#", "##", "###")):
    blocks = []
    n = len(lines)
    stop_prefixes = tuple(lvl + " " for lvl in stop_levels)
    while i < n:
        line = lines[i]
        stripped = line.strip()
        if stripped == "":
            i += 1
            continue
        if line.startswith(stop_prefixes):
            break
        if stripped == "---":
            i += 1
            continue
        if stripped.startswith("```mermaid"):
            i += 1
            code_lines = []
            while i < n and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append({"type": "mermaid", "code": "\n".join(code_lines)})
            continue
        if stripped.startswith("```"):
            # bloque de código genérico: se preserva como texto monoespaciado
            i += 1
            code_lines = []
            while i < n and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append({"type": "code", "code": "\n".join(code_lines)})
            continue
        if stripped.startswith("$$"):
            if stripped.endswith("$$") and len(stripped) > 4:
                blocks.append({"type": "displaymath", "tex": stripped[2:-2]})
                i += 1
                continue
            math_lines = [stripped[2:]]
            i += 1
            while i < n and "$$" not in lines[i]:
                math_lines.append(lines[i])
                i += 1
            if i < n:
                math_lines.append(lines[i].split("$$")[0])
                i += 1
            blocks.append({"type": "displaymath", "tex": "\n".join(math_lines).strip()})
            continue
        if stripped.startswith("|"):
            table_lines = []
            while i < n and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            blocks.append({"type": "table", "rows": parse_table(table_lines)})
            continue
        if stripped.startswith(">"):
            bq_lines = []
            while i < n and lines[i].strip().startswith(">"):
                content = lines[i].strip()[1:]
                if content.startswith(" "):
                    content = content[1:]
                bq_lines.append(content.strip())
                i += 1
            blocks.append({"type": "blockquote", "blocks": parse_blockquote(bq_lines)})
            continue
        if re.match(r"^[-*]\s+", stripped):
            items = []
            base_indent = None
            while i < n and re.match(r"^[-*]\s+", lines[i].strip()):
                raw = lines[i]
                indent = len(raw) - len(raw.lstrip(" "))
                if base_indent is None:
                    base_indent = indent
                level = max(0, (indent - base_indent) // 2)
                items.append((re.sub(r"^[-*]\s+", "", raw.strip()), level))
                i += 1
            blocks.append({"type": "bullets", "items": items})
            continue
        if re.match(r"^\d+\.\s+", stripped):
            items = []
            while i < n and re.match(r"^\d+\.\s+", lines[i].strip()):
                items.append((re.sub(r"^\d+\.\s+", "", lines[i].strip()), 0))
                i += 1
            blocks.append({"type": "numbered", "items": items})
            continue
        if re.match(r"^\*[^*].*\*$", stripped) and stripped.count("*") == 2:
            blocks.append({"type": "citation", "text": stripped[1:-1]})
            i += 1
            continue
        para_lines = [line.strip()]
        i += 1
        while i < n:
            nxt = lines[i]
            nxt_s = nxt.strip()
            if nxt_s == "" or nxt.startswith(stop_prefixes) or nxt_s.startswith("|") or \
               nxt_s.startswith(">") or re.match(r"^[-*]\s+", nxt_s) or \
               re.match(r"^\d+\.\s+", nxt_s) or nxt_s.startswith("```") or \
               nxt_s.startswith("$$") or nxt_s == "---":
                break
            para_lines.append(nxt_s)
            i += 1
        blocks.append({"type": "para", "text": " ".join(para_lines)})
    return blocks, i


def parse_md(path):
    with io.open(path, encoding="utf-8") as f:
        text = f.read()
    lines = text.split("\n")
    n = len(lines)
    i = 0
    while i < n and not lines[i].startswith("# "):
        i += 1
    full_title = lines[i][2:].strip()
    title = full_title.split("·", 1)[-1].strip() if "·" in full_title else full_title
    i += 1
    doc = {"full_title": full_title, "title": title, "subtitle": "",
           "objetivo": "", "contenido": [], "sections": [], "fuentes": [],
           "cierre_title": "Cierre de la unidad", "cierre": []}
    while i < n and lines[i].strip() == "":
        i += 1
    if i < n and lines[i].startswith("**"):
        doc["subtitle"] = re.sub(r"\*\*", "", lines[i]).strip()
        i += 1

    while i < n:
        line = lines[i]
        if line.startswith("## "):
            heading = line[3:].strip()
            if heading.startswith("Objetivo de la unidad"):
                i += 1
                blocks, i = parse_blocks_until_heading(lines, i, ("#", "##"))
                doc["objetivo"] = " ".join(
                    plain_text_of(b["text"]) for b in blocks if b["type"] == "para"
                )
                continue
            if heading.startswith("Contenido"):
                i += 1
                while i < n and lines[i].strip() == "":
                    i += 1
                table_lines = []
                while i < n and lines[i].strip().startswith("|"):
                    table_lines.append(lines[i])
                    i += 1
                doc["contenido"] = parse_table(table_lines)
                # consumir bloques intermedios (notas de práctica) sin usarlos
                _, i = parse_blocks_until_heading(lines, i, ("#", "##", "###"))
                continue
            if heading.startswith("Fuentes y referencias"):
                i += 1
                blocks, i = parse_blocks_until_heading(lines, i, ("#", "##"))
                doc["fuentes"] = blocks
                continue
            if heading.startswith("Cierre de la unidad"):
                doc["cierre_title"] = heading
                i += 1
                blocks, i = parse_blocks_until_heading(lines, i, ("#", "##"))
                doc["cierre"] = blocks
                continue
            # cualquier otro ## (p.ej. "Apéndice: ...") se trata como sección de contenido
            i += 1
            skip_slide, i = consume_skip_marker(lines, i)
            blocks, i = parse_blocks_until_heading(lines, i, ("#", "##"))
            doc["sections"].append({"title": heading, "number": None, "blocks": blocks,
                                     "skip_slide": skip_slide})
            continue
        if line.startswith("### "):
            heading = line[4:].strip()
            m = re.match(r"(\d+)\.\s*(.*)", heading)
            num = int(m.group(1)) if m else None
            title_s = m.group(2) if m else heading
            i += 1
            skip_slide, i = consume_skip_marker(lines, i)
            blocks, i = parse_blocks_until_heading(lines, i, ("#", "##", "###"))
            doc["sections"].append({"title": title_s, "number": num, "blocks": blocks,
                                     "skip_slide": skip_slide})
            continue
        i += 1
    return doc


# --------------------------------------------------------------------------
# Renderizado de fórmulas (matplotlib mathtext) y diagramas mermaid
# --------------------------------------------------------------------------

def render_formula_png(tex, out_path, fontsize=24):
    tex = tex.replace(r"\dfrac", r"\frac")
    fig = plt.figure(figsize=(0.1, 0.1))
    txt = fig.text(0, 0, f"${tex}$", fontsize=fontsize)
    fig.canvas.draw()
    bbox = txt.get_window_extent()
    w_in = bbox.width / fig.dpi + 0.15
    h_in = bbox.height / fig.dpi + 0.12
    plt.close(fig)
    fig = plt.figure(figsize=(w_in, h_in))
    fig.text(0.5, 0.5, f"${tex}$", fontsize=fontsize, ha="center", va="center",
              color="#1E2761")
    fig.savefig(out_path, dpi=220, transparent=True)
    plt.close(fig)
    return w_in, h_in


MERMAID_EDGE_RE = re.compile(
    r"^(\w+)\s*(-\.->|-->|-\.-)\s*(?:\|([^|]*)\|)?\s*(\w+)"
)
MERMAID_CIRCLE_RE = re.compile(r'^(\w+)\(\(\s*["\']?(.*?)["\']?\s*\)\)')
MERMAID_BOX_RE = re.compile(r'^(\w+)\[\s*["\']?(.*?)["\']?\s*\]')


def parse_mermaid_graph(code):
    nodes, order, edges = {}, [], []

    def ensure(nid, label=None, shape="box"):
        if nid not in nodes:
            nodes[nid] = {"label": label if label else nid, "shape": shape}
            order.append(nid)
        elif label:
            nodes[nid]["label"] = label
            nodes[nid]["shape"] = shape

    for raw in code.split("\n"):
        line = raw.strip().rstrip(";")
        if not line or line.startswith("graph") or line.startswith("flowchart"):
            continue
        m = MERMAID_EDGE_RE.match(line)
        if m:
            src, style, label, dst = m.groups()
            ensure(src)
            ensure(dst)
            edges.append((src, dst, style, label))
            continue
        m = MERMAID_CIRCLE_RE.match(line)
        if m:
            ensure(m.group(1), m.group(2), "circle")
            continue
        m = MERMAID_BOX_RE.match(line)
        if m:
            ensure(m.group(1), m.group(2), "box")
            continue
    return nodes, order, edges


def render_mermaid_png(code, out_path):
    import networkx as nx
    from collections import defaultdict

    nodes, order, edges = parse_mermaid_graph(code)
    G = nx.DiGraph()
    G.add_nodes_from(order)
    # Layering usa toda arista con flecha (sólida o punteada): ambas expresan
    # una dependencia dirigida. Solo "-.-" (coordinación sin flecha) se excluye,
    # porque no implica que el destino deba quedar en un nivel inferior.
    directed = [(s, d) for s, d, style, _ in edges if style in ("-->", "-.->")]
    G.add_edges_from(directed)
    layer = {n: 0 for n in order}
    try:
        topo = list(nx.topological_sort(G))
    except Exception:
        topo = order
    for n in topo:
        preds = list(G.predecessors(n))
        if preds:
            layer[n] = max(layer[p] for p in preds) + 1
    by_layer = defaultdict(list)
    for n in order:
        by_layer[layer[n]].append(n)
    n_layers = max(by_layer.keys()) + 1 if by_layer else 1

    fig_w, fig_h = 11.5, max(3.0, n_layers * 1.5)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, n_layers)
    ax.invert_yaxis()
    ax.axis("off")

    pos = {}
    for lvl, ns in by_layer.items():
        count = len(ns)
        for idx, nid in enumerate(ns):
            x = (idx + 1) * 10.0 / (count + 1)
            y = lvl + 0.5
            pos[nid] = (x, y)

    box_w, box_h = 1.55, 0.55
    for nid in order:
        x, y = pos[nid]
        label = nodes[nid]["label"]
        if nodes[nid]["shape"] == "circle":
            patch = mpatches.Circle((x, y), 0.42, facecolor="#1E2761",
                                     edgecolor="#C9A227", linewidth=1.5, zorder=2)
        else:
            patch = mpatches.FancyBboxPatch(
                (x - box_w / 2, y - box_h / 2), box_w, box_h,
                boxstyle="round,pad=0.02,rounding_size=0.08",
                facecolor="#1E2761", edgecolor="#1E2761", zorder=2)
        ax.add_patch(patch)
        ax.text(x, y, label, ha="center", va="center", color="white",
                 fontsize=11, fontweight="bold", zorder=3, wrap=True)

    for src, dst, style, label in edges:
        x1, y1 = pos[src]
        x2, y2 = pos[dst]
        dotted = style in ("-.-", "-.->")
        arrow = style in ("-->", "-.->")
        arrowstyle = "-|>" if arrow else "-"
        con = mpatches.FancyArrowPatch(
            (x1, y1), (x2, y2), arrowstyle=arrowstyle,
            linestyle=":" if dotted else "-", color="#5B6482",
            mutation_scale=14, linewidth=1.4, shrinkA=22, shrinkB=22, zorder=1)
        ax.add_patch(con)
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            ax.text(mx, my, label, ha="center", va="center", fontsize=9,
                     color="#27314D", style="italic", zorder=4,
                     bbox=dict(boxstyle="round,pad=0.15", fc="white", ec="none"))

    fig.tight_layout(pad=0.3)
    fig.savefig(out_path, dpi=200, transparent=True)
    plt.close(fig)
    from PIL import Image
    with Image.open(out_path) as im:
        return im.width, im.height


# --------------------------------------------------------------------------
# Helpers de bajo nivel para construir shapes de texto
# --------------------------------------------------------------------------

def _apply_run(run, text, bold, italic, code, size_pt, color, font=FONT_BODY):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    run.font.name = FONT_CODE if code else font
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, size_pt=14, color=BODY_COLOR,
                 bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT_BODY,
                 anchor=None, word_wrap=True):
    box = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    return box


def set_paragraph_plain(paragraph, text, size_pt, color, bold=False, italic=False,
                         align=PP_ALIGN.LEFT, font=FONT_BODY):
    paragraph.alignment = align
    run = paragraph.add_run()
    _apply_run(run, text, bold, italic, False, size_pt, color, font)


def set_paragraph_runs(paragraph, runs, size_pt, color, align=PP_ALIGN.LEFT,
                        bold_color=None, bullet_prefix=None):
    paragraph.alignment = align
    if bullet_prefix:
        r = paragraph.add_run()
        _apply_run(r, bullet_prefix, True, False, False, size_pt, NAVY)
    for (t, b, i, c) in runs:
        r = paragraph.add_run()
        rc = bold_color if (b and bold_color) else color
        _apply_run(r, t, b, i, c, size_pt, rc)


def add_eyebrow_title(slide, eyebrow, title, page_num=None, title_size=32):
    if eyebrow:
        box = add_textbox(slide, MARGIN_L, EYEBROW_TOP, CONTENT_W, EYEBROW_H)
        p = box.text_frame.paragraphs[0]
        set_paragraph_plain(p, eyebrow.upper(), 12, GOLD, bold=True)
        for run in p.runs:
            run.font._rPr.set("spc", "150")
    lines = count_wrapped_lines(title, emu_to_pt72(CONTENT_W), title_size, bold=True, font=FONT_TITLE)
    title_h = max(700000, lines * (title_size * 1.22) * EMU_PER_PT)
    box = add_textbox(slide, MARGIN_L, TITLE_TOP, CONTENT_W, title_h)
    p = box.text_frame.paragraphs[0]
    set_paragraph_plain(p, title, title_size, NAVY, bold=True, font=FONT_TITLE)
    if page_num is not None:
        pbox = add_textbox(slide, SLIDE_W - MARGIN_R - 548640, PAGENUM_TOP, 548640, 274320,
                            align=PP_ALIGN.RIGHT)
        p2 = pbox.text_frame.paragraphs[0]
        set_paragraph_plain(p2, f"{page_num:02d}", 10, LGRAY, align=PP_ALIGN.RIGHT)
    return TITLE_TOP + title_h


def add_footer_citation(slide, text):
    box = add_textbox(slide, MARGIN_L, FOOTER_TOP, 9418320, 274320)
    p = box.text_frame.paragraphs[0]
    set_paragraph_plain(p, text, 9.5, GRAY, italic=True)


# --------------------------------------------------------------------------
# Bloques de contenido: estimar altura y dibujar
# --------------------------------------------------------------------------

LINE_SP = 1.28
GAP = 90000


def estimate_para_height(text_or_runs, width_emu, size_pt=14, bold=False):
    if isinstance(text_or_runs, str):
        plain = text_or_runs
    else:
        plain = "".join(t for t, b, i, c in text_or_runs)
    lines = count_wrapped_lines(plain, emu_to_pt72(width_emu), size_pt, bold=bold)
    return int(lines * size_pt * LINE_SP * EMU_PER_PT + 30000)


def draw_para(slide, runs, left, top, width, size_pt=14, color=BODY_COLOR,
              align=PP_ALIGN.LEFT, font=FONT_BODY):
    plain = "".join(t for t, b, i, c in runs)
    lines = count_wrapped_lines(plain, emu_to_pt72(width), size_pt)
    h = int(lines * size_pt * LINE_SP * EMU_PER_PT + 30000)
    box = add_textbox(slide, left, top, width, h)
    p = box.text_frame.paragraphs[0]
    set_paragraph_runs(p, runs, size_pt, color, align=align)
    for r in box.text_frame.paragraphs[0].runs:
        r.font.name = font if not r.font.name else r.font.name
    return h


BULLET_MARKER_W = 260000
BULLET_INDENT_STEP = 260000


def _item_text_level(it):
    return it if isinstance(it, tuple) else (it, 0)


def estimate_bullets_height(items, width_emu, size_pt=13.5):
    total = 0
    for it in items:
        text, level = _item_text_level(it)
        runs = parse_inline(text)
        plain = "".join(t for t, b, i, c in runs)
        avail = width_emu - BULLET_MARKER_W - level * BULLET_INDENT_STEP
        lines = count_wrapped_lines(plain, emu_to_pt72(avail), size_pt)
        total += int(lines * size_pt * LINE_SP * EMU_PER_PT + 55000)
    return total


def draw_bullets(slide, items, left, top, width, size_pt=13.5, numbered=False,
                  color=BODY_COLOR, marker_color=GOLD):
    y = top
    for idx, it in enumerate(items):
        text, level = _item_text_level(it)
        indent = level * BULLET_INDENT_STEP
        runs = parse_inline(text)
        plain = "".join(t for t, b, i, c in runs)
        avail = width - BULLET_MARKER_W - indent
        lines = count_wrapped_lines(plain, emu_to_pt72(avail), size_pt)
        h = int(lines * size_pt * LINE_SP * EMU_PER_PT + 55000)
        mbox = add_textbox(slide, left + indent, y, BULLET_MARKER_W, h)
        mp = mbox.text_frame.paragraphs[0]
        marker = f"{idx + 1}." if numbered else ("◦" if level else "•")
        set_paragraph_plain(mp, marker, size_pt, marker_color, bold=True)
        tbox = add_textbox(slide, left + indent + BULLET_MARKER_W, y, avail, h)
        tp = tbox.text_frame.paragraphs[0]
        set_paragraph_runs(tp, runs, size_pt, color)
        y += h
    return y - top


def plain_cell(text):
    return plain_text_of(text)


def compute_col_widths(rows, total_width):
    ncols = len(rows[0])
    maxlen = [1] * ncols
    for row in rows:
        for j in range(ncols):
            cell = row[j] if j < len(row) else ""
            maxlen[j] = max(maxlen[j], len(plain_cell(cell)))
    total = sum(maxlen)
    # Piso de 11%: por debajo de eso, una sola palabra larga en la columna
    # (p. ej. "Instrumentos" en una tabla comparativa de dos columnas anchas)
    # no cabe y PowerPoint la corta a media palabra en vez de hacer wrap.
    widths = [max(int(total_width * ml / total), int(total_width * 0.11)) for ml in maxlen]
    scale = total_width / sum(widths)
    widths = [int(w * scale) for w in widths]
    widths[-1] += total_width - sum(widths)
    return widths


def table_font_size(rows):
    ncols = len(rows[0])
    if ncols <= 3:
        return 12.5
    if ncols == 4:
        return 11.5
    if ncols == 5:
        return 10.5
    return 9.5


def row_height(row, col_widths, size_pt, header=False):
    max_lines = 1
    for j, w in enumerate(col_widths):
        cell = row[j] if j < len(row) else ""
        txt = plain_cell(cell)
        avail_pt = emu_to_pt72(w) - 10
        lines = count_wrapped_lines(txt, avail_pt, size_pt, bold=header)
        max_lines = max(max_lines, lines)
    return int(max_lines * size_pt * 1.32 * EMU_PER_PT + 118000)


def estimate_table_height(rows, col_widths, size_pt):
    h = row_height(rows[0], col_widths, size_pt + 0.5, header=True)
    for r in rows[1:]:
        h += row_height(r, col_widths, size_pt)
    return h


def split_table_rows(rows, col_widths, size_pt, max_height):
    header = rows[0]
    header_h = row_height(header, col_widths, size_pt + 0.5, header=True)
    chunks = []
    cur = []
    cur_h = header_h
    for row in rows[1:]:
        h = row_height(row, col_widths, size_pt)
        if cur and cur_h + h > max_height:
            chunks.append(cur)
            cur = []
            cur_h = header_h
        cur.append(row)
        cur_h += h
    if cur:
        chunks.append(cur)
    if not chunks:
        chunks = [[]]
    return chunks


def draw_table(slide, header, body_rows, col_widths, left, top, size_pt):
    n_rows = len(body_rows) + 1
    n_cols = len(header)
    total_w = sum(col_widths)
    total_h = estimate_table_height([header] + body_rows, col_widths, size_pt)
    gframe = slide.shapes.add_table(n_rows, n_cols, Emu(int(left)), Emu(int(top)),
                                     Emu(int(total_w)), Emu(int(total_h)))
    table = gframe.table
    table.first_row = False
    table.horz_banding = False
    tbl_el = table._tbl
    style_el = tbl_el.find(qn("a:tblPr")).find(qn("a:tableStyleId"))
    if style_el is not None:
        style_el.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    for j, w in enumerate(col_widths):
        table.columns[j].width = Emu(int(w))
    hdr_h = row_height(header, col_widths, size_pt + 0.5, header=True)
    table.rows[0].height = Emu(int(hdr_h))
    for j, cell_text in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEAD_BG
        cell.margin_left = cell.margin_right = Emu(54864)
        cell.margin_top = cell.margin_bottom = Emu(27432)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        runs = parse_inline(cell_text)
        set_paragraph_runs(p, runs, size_pt, WHITE)
        for r in p.runs:
            r.font.bold = True
    for i, row in enumerate(body_rows, start=1):
        rh = row_height(row, col_widths, size_pt)
        table.rows[i].height = Emu(int(rh))
        for j in range(n_cols):
            cell_text = row[j] if j < len(row) else ""
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT if i % 2 == 0 else WHITE
            cell.margin_left = cell.margin_right = Emu(54864)
            cell.margin_top = cell.margin_bottom = Emu(27432)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            runs = parse_inline(cell_text)
            set_paragraph_runs(p, runs, size_pt, BODY_COLOR)
    return total_h


def estimate_blockquote_height(bq_blocks, width_emu):
    inner_pad = 180000
    inner_w = width_emu - 2 * inner_pad
    h = 2 * 140000
    for b in bq_blocks:
        if b["type"] == "para":
            h += estimate_para_height(b["text"], inner_w, 13) + 40000
        elif b["type"] == "bullets":
            h += estimate_bullets_height(b["items"], inner_w, 13) + 40000
    return h


def draw_blockquote(slide, bq_blocks, left, top, width):
    inner_pad = 180000
    inner_w = width - 2 * inner_pad
    h = estimate_blockquote_height(bq_blocks, width)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(left)), Emu(int(top)),
                                   Emu(int(width)), Emu(int(h)))
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.fill.background()
    rect.shadow.inherit = False
    try:
        rect.adjustments[0] = 0.05
    except Exception:
        pass
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)),
                                  Emu(45720), Emu(int(h)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    bar.shadow.inherit = False
    y = top + 140000
    for b in bq_blocks:
        if b["type"] == "para":
            runs = parse_inline(b["text"])
            hh = draw_para(slide, runs, left + inner_pad, y, inner_w, size_pt=13, color=NAVY)
            y += hh + 40000
        elif b["type"] == "bullets":
            hh = draw_bullets(slide, b["items"], left + inner_pad, y, inner_w, size_pt=13,
                               color=NAVY, marker_color=GOLD)
            y += hh + 40000
    return h


def draw_image_block(slide, img_path, img_w, img_h, left, top, max_width):
    scale = min(1.0, max_width / img_w) if img_w else 1.0
    w = int(img_w * scale)
    h = int(img_h * scale)
    x = int(left + (max_width - w) / 2)
    slide.shapes.add_picture(img_path, Emu(x), Emu(int(top)), Emu(w), Emu(h))
    return h


def estimate_block_height(b, width_emu):
    t = b["type"]
    if t == "para":
        return estimate_para_height(b["text"], width_emu, 14) + GAP
    if t in ("bullets", "numbered"):
        return estimate_bullets_height(b["items"], width_emu, 13.5) + GAP
    if t == "table":
        col_widths = compute_col_widths(b["rows"], width_emu)
        size_pt = table_font_size(b["rows"])
        return estimate_table_height(b["rows"], col_widths, size_pt) + GAP
    if t == "blockquote":
        return estimate_blockquote_height(b["blocks"], width_emu) + GAP
    if t == "citation":
        return estimate_para_height(b["text"], width_emu, 10.5, bold=False) + GAP
    if t in ("mermaid", "displaymath"):
        return b.get("img_h_emu", 1500000) + GAP
    if t == "code":
        n_lines = len(b["code"].split("\n"))
        return int(n_lines * 12 * 1.3 * EMU_PER_PT + 60000) + GAP
    return GAP


def draw_block(slide, b, left, top, width_emu):
    t = b["type"]
    if t == "para":
        runs = parse_inline(b["text"])
        return draw_para(slide, runs, left, top, width_emu, size_pt=14, color=BODY_COLOR)
    if t in ("bullets", "numbered"):
        return draw_bullets(slide, b["items"], left, top, width_emu, size_pt=13.5,
                             numbered=(t == "numbered"))
    if t == "table":
        col_widths = compute_col_widths(b["rows"], width_emu)
        size_pt = table_font_size(b["rows"])
        return draw_table(slide, b["rows"][0], b["rows"][1:], col_widths, left, top, size_pt)
    if t == "blockquote":
        return draw_blockquote(slide, b["blocks"], left, top, width_emu)
    if t == "citation":
        return draw_para(slide, [(b["text"], False, True, False)], left, top, width_emu,
                          size_pt=10.5, color=GRAY)
    if t in ("mermaid", "displaymath"):
        return draw_image_block(slide, b["img_path"], b["img_w_emu"], b["img_h_emu"],
                                 left, top, width_emu)
    if t == "code":
        n_lines = len(b["code"].split("\n"))
        h = int(n_lines * 12 * 1.3 * EMU_PER_PT + 60000)
        box = add_textbox(slide, left, top, width_emu, h)
        for j, ln in enumerate(b["code"].split("\n")):
            p = box.text_frame.paragraphs[0] if j == 0 else box.text_frame.add_paragraph()
            set_paragraph_plain(p, ln, 11, BODY_COLOR, font=FONT_CODE)
        return h
    return 0


# --------------------------------------------------------------------------
# Preprocesamiento: renderiza mermaid / display-math a PNG una sola vez
# --------------------------------------------------------------------------

def prepare_media(blocks, tmpdir, prefix):
    for k, b in enumerate(blocks):
        if b["type"] == "mermaid":
            out = os.path.join(tmpdir, f"{prefix}_mermaid_{k}.png")
            w_px, h_px = render_mermaid_png(b["code"], out)
            max_w = CONTENT_W
            max_h = 2600000
            scale = min(max_w / w_px, max_h / h_px)
            b["img_path"] = out
            b["img_w_emu"] = w_px * scale
            b["img_h_emu"] = h_px * scale
        elif b["type"] == "displaymath":
            out = os.path.join(tmpdir, f"{prefix}_math_{k}.png")
            w_in, h_in = render_formula_png(b["tex"], out)
            w_emu = w_in * EMU_PER_IN
            h_emu = h_in * EMU_PER_IN
            max_w = CONTENT_W * 0.85
            if w_emu > max_w:
                scale = max_w / w_emu
                w_emu *= scale
                h_emu *= scale
            b["img_path"] = out
            b["img_w_emu"] = w_emu
            b["img_h_emu"] = h_emu
        elif b["type"] == "blockquote":
            prepare_media(b["blocks"], tmpdir, prefix + f"_bq{k}")


# --------------------------------------------------------------------------
# Construcción del deck
# --------------------------------------------------------------------------

CONTENT_START_DEFAULT = 1850000


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self.blank_layout = self.prs.slide_layouts[6]
        self.page = 1

    def new_slide(self):
        return self.prs.slides.add_slide(self.blank_layout)

    def next_page(self):
        self.page += 1
        return self.page


def flow_blocks(deck, eyebrow, base_title, blocks, allow_subtitle=True):
    """Genera 1+ slides para `blocks`, paginando automáticamente.
    Devuelve el número de página de la primera slide creada."""
    blocks = list(blocks)
    first_page = None
    cont = False
    if not blocks:
        blocks = [{"type": "para", "text": ""}]
    idx = 0
    while idx < len(blocks):
        slide = deck.new_slide()
        page = deck.next_page()
        if first_page is None:
            first_page = page
        title = base_title + (" (cont.)" if cont else "")
        title_bottom = add_eyebrow_title(slide, eyebrow, title, page_num=page)
        y = title_bottom + 90000
        if allow_subtitle and not cont and blocks[idx]["type"] == "para" and \
                len(plain_text_of(blocks[idx]["text"])) <= 200:
            runs = parse_inline(blocks[idx]["text"])
            h = draw_para(slide, runs, MARGIN_L, y, CONTENT_W, size_pt=15.5, color=GRAY)
            y += h + 130000
            idx += 1
        placed_any = False
        while idx < len(blocks):
            b = blocks[idx]
            h = estimate_block_height(b, CONTENT_W)
            bottom_limit = BODY_BOTTOM_MAX + (260000 if b["type"] == "citation" else 0)
            avail = bottom_limit - y
            if b["type"] == "para" and idx + 1 < len(blocks) and placed_any and \
                    len(plain_text_of(b["text"])) < 160 and \
                    plain_text_of(b["text"]).rstrip().endswith(":"):
                nxt_h = estimate_block_height(blocks[idx + 1], CONTENT_W)
                if h + nxt_h > avail:
                    break
            if h > avail and placed_any:
                break
            if h > avail and b["type"] == "table":
                col_widths = compute_col_widths(b["rows"], CONTENT_W)
                size_pt = table_font_size(b["rows"])
                chunks = split_table_rows(b["rows"], col_widths, size_pt, max(avail - GAP, 400000))
                header = b["rows"][0]
                first_chunk = chunks[0]
                hh = draw_table(slide, header, first_chunk, col_widths, MARGIN_L, y, size_pt)
                y += hh + GAP
                placed_any = True
                rest = chunks[1:]
                if rest:
                    blocks[idx] = {"type": "table",
                                    "rows": [header] + [r for c in rest for r in c]}
                else:
                    idx += 1
                break
            hh = draw_block(slide, b, MARGIN_L, y, CONTENT_W)
            y += hh + GAP
            idx += 1
            placed_any = True
        cont = True
    return first_page


def add_cover_slide(deck, doc):
    slide = deck.new_slide()
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9600000), Emu(-1400000),
                                     Emu(4200000), Emu(4200000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    circle.shadow.inherit = False
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(10200000), Emu(4300000),
                                      Emu(3200000), Emu(3200000))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = NAVY
    circle2.line.fill.background()
    circle2.shadow.inherit = False

    ebox = add_textbox(slide, MARGIN_L, 1650000, CONTENT_W, EYEBROW_H)
    ep = ebox.text_frame.paragraphs[0]
    set_paragraph_plain(ep, "MERCADOS DE DEUDA Y CAPITALES", 13, GOLD, bold=True)

    title_lines = count_wrapped_lines(doc["title"], emu_to_pt72(9200000), 40, bold=True, font=FONT_TITLE)
    tbox = add_textbox(slide, MARGIN_L, 2050000, 9200000, title_lines * 620000)
    tp = tbox.text_frame.paragraphs[0]
    set_paragraph_plain(tp, doc["title"], 40, NAVY, bold=True, font=FONT_TITLE)

    sub_top = 2050000 + title_lines * 620000 + 120000
    sbox = add_textbox(slide, MARGIN_L, sub_top, 9200000, 400000)
    sp = sbox.text_frame.paragraphs[0]
    set_paragraph_plain(sp, doc["subtitle"], 13, GRAY, italic=True)
    return slide


def add_objetivo_slide(deck, doc, page):
    slide = deck.new_slide()
    add_eyebrow_title(slide, doc["title"], "Objetivo de la unidad", page_num=page)
    box_top = 2743200
    runs = parse_inline(doc["objetivo"])
    lines = count_wrapped_lines(doc["objetivo"], emu_to_pt72(10058400), 20, bold=False)
    box_h = max(1200000, int(lines * 20 * 1.5 * EMU_PER_PT) + 500000)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(MARGIN_L), Emu(box_top),
                                   Emu(10881360), Emu(box_h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()
    card.shadow.inherit = False
    try:
        card.adjustments[0] = 0.04
    except Exception:
        pass
    tbox = add_textbox(slide, MARGIN_L + 548640, box_top + 200000, 10058400, box_h - 400000,
                        anchor=MSO_ANCHOR.MIDDLE)
    p = tbox.text_frame.paragraphs[0]
    set_paragraph_runs(p, runs, 20, BODY_COLOR)
    for r in p.runs:
        r.font.name = FONT_TITLE
    return slide


def add_cierre_slides(deck, title, items):
    """items: lista de strings markdown (cada uno un bullet de cierre)."""
    blocks = [{"type": "para_bold_first", "text": it} for it in items]
    idx = 0
    while idx < len(blocks):
        slide = deck.new_slide()
        title_bottom = add_eyebrow_title(slide, "", f"Cierre · {title}")
        y = title_bottom + 180000
        placed = 0
        while idx < len(blocks):
            runs = parse_inline(blocks[idx]["text"])
            plain = "".join(t for t, b, i, c in runs)
            lines = count_wrapped_lines(plain, emu_to_pt72(10332720), 16)
            h = int(lines * 16 * 1.35 * EMU_PER_PT + 60000)
            if y + h > BODY_BOTTOM_MAX and placed > 0:
                break
            box = add_textbox(slide, 822960, y, 10332720, h)
            p = box.text_frame.paragraphs[0]
            set_paragraph_runs(p, runs, 16, NAVY, bold_color=NAVY)
            y += h + 200000
            idx += 1
            placed += 1
    return


def build_deck(doc, out_pptx):
    tmpdir = tempfile.mkdtemp(prefix="mdpptx_")
    for s in doc["sections"]:
        prepare_media(s["blocks"], tmpdir, f"sec{s['title'][:10]}")

    deck = Deck()
    add_cover_slide(deck, doc)
    add_objetivo_slide(deck, doc, deck.next_page())

    # Reservar slide de TOC: se crea vacía y se llena al final para poder
    # calcular los números de página reales de cada sección.
    toc_slide = deck.new_slide()
    deck.next_page()

    contenido_rows = doc["contenido"][1:] if doc["contenido"] else []
    entries = []
    for row_idx, row in enumerate(contenido_rows):
        section = doc["sections"][row_idx] if row_idx < len(doc["sections"]) else None
        if section is not None and section.get("skip_slide"):
            continue
        tema = row[1] if len(row) > 1 else ""
        quecubre = row[2] if len(row) > 2 else ""
        entries.append([None, tema, quecubre, None])
    for entry_idx, e in enumerate(entries):
        e[0] = to_roman(entry_idx + 1)

    entry_cursor = 0
    for row_idx, section in enumerate(doc["sections"]):
        if section.get("skip_slide"):
            continue
        eyebrow = doc["title"]
        heading = section["title"]
        first_page = flow_blocks(deck, eyebrow, heading, section["blocks"])
        if row_idx < len(contenido_rows) and entry_cursor < len(entries):
            entries[entry_cursor][3] = first_page
            entry_cursor += 1

    for e in entries:
        if e[3] is None:
            e[3] = deck.page

    cierre_items = []
    for blk in doc["cierre"]:
        if blk["type"] == "bullets":
            cierre_items.extend(text for text, level in blk["items"])
        elif blk["type"] == "para":
            cierre_items.append(blk["text"])
    add_cierre_slides(deck, doc["title"], cierre_items)

    fuentes_items = []
    for blk in doc["fuentes"]:
        if blk["type"] == "bullets":
            fuentes_items.extend(text for text, level in blk["items"])
        elif blk["type"] == "para":
            fuentes_items.append(blk["text"])
    fuentes_page = flow_blocks(deck, "PREPARACIÓN DOCENTE", "Fuentes y referencias recomendadas",
                                [{"type": "bullets", "items": fuentes_items}],
                                allow_subtitle=False)

    # Rellenar la slide de TOC reservada (índice 2, ya en su posición correcta)
    _fill_toc_slide(toc_slide, entries)

    deck.prs.save(out_pptx)
    return out_pptx


def _fill_toc_slide(slide, entries):
    add_eyebrow_title(slide, "", "Tabla de contenido", page_num=3)
    y = 1900000
    row_h = int((BODY_BOTTOM_MAX - y) / max(len(entries), 1))
    row_h = min(row_h, 620000)
    for roman, tema, quecubre, pref in entries:
        rbox = add_textbox(slide, MARGIN_L, y, 700000, 400000)
        rp = rbox.text_frame.paragraphs[0]
        set_paragraph_plain(rp, roman + ".", 15, GOLD, bold=True, font=FONT_TITLE)
        tbox = add_textbox(slide, MARGIN_L + 750000, y, 8600000, 400000)
        tp = tbox.text_frame.paragraphs[0]
        set_paragraph_plain(tp, tema, 15, NAVY, bold=True, font=FONT_TITLE)
        if quecubre:
            qbox = add_textbox(slide, MARGIN_L + 750000, y + 330000, 8600000, 300000)
            qp = qbox.text_frame.paragraphs[0]
            set_paragraph_plain(qp, quecubre, 10.5, GRAY)
        pbox = add_textbox(slide, SLIDE_W - MARGIN_R - 1200000, y, 1200000, 400000,
                            align=PP_ALIGN.RIGHT)
        pp = pbox.text_frame.paragraphs[0]
        set_paragraph_plain(pp, f"PÁG. {pref:02d}", 11, GRAY, align=PP_ALIGN.RIGHT)
        y += row_h


# --------------------------------------------------------------------------
# Exportar a PDF vía PowerPoint (Windows)
# --------------------------------------------------------------------------

def export_pdf(pptx_path, pdf_path):
    import win32com.client
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = os.path.abspath(pdf_path)
    app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        pres = app.Presentations.Open(pptx_path, True, False, False)
        pres.SaveAs(pdf_path, 32)
        pres.Close()
    finally:
        app.Quit()
    return pdf_path


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    md_path = sys.argv[1]
    out_pptx = sys.argv[2] if len(sys.argv) > 2 else os.path.splitext(md_path)[0] + ".pptx"
    doc = parse_md(md_path)
    build_deck(doc, out_pptx)
    print(f"pptx generado: {out_pptx}")
    if "--no-pdf" not in sys.argv:
        pdf_path = os.path.splitext(out_pptx)[0] + ".pdf"
        export_pdf(out_pptx, pdf_path)
        print(f"pdf generado: {pdf_path}")


if __name__ == "__main__":
    main()
